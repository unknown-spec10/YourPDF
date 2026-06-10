import os
import uuid
import subprocess
import logging
import zipfile
import sys
import shutil
from pypdf import PdfReader, PdfWriter
from pdf2image import convert_from_path
import img2pdf
from pdf2docx import Converter
from app.celery_app import celery
from app.s3 import store_processed_file, is_s3_configured
from PIL import Image, ImageDraw, ImageFont

logger = logging.getLogger(__name__)

@celery.task
def delete_local_file_task(file_path: str):
    """
    Background task to clean up local output files after the expiry period (15 minutes).
    """
    if os.path.exists(file_path):
        try:
            os.remove(file_path)
            logger.info(f"Successfully deleted local fallback file: {file_path}")
        except Exception as e:
            logger.error(f"Failed to delete local fallback file {file_path}: {e}")

@celery.task(bind=True)
def compress_pdf_task(self, file_path: str, quality: str, original_filename: str = None):
    """
    Runs Ghostscript compression on a PDF file, uploads to S3, and returns results.
    """
    logger.info(f"Starting compression task for {file_path} with quality {quality}")
    
    # 20% - Initialize
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Initializing compression..."})
    
    # Resolve relative paths relative to backend root
    file_path = file_path.replace('\\', '/')
    if not os.path.isabs(file_path):
        app_dir = os.path.dirname(os.path.abspath(__file__))
        backend_root = os.path.dirname(app_dir)
        file_path = os.path.abspath(os.path.join(backend_root, file_path))
        
    if not os.path.exists(file_path):
        logger.error(f"Input file not found: {file_path}")
        raise FileNotFoundError(f"Input file not found: {file_path}")

    # Map quality presets (representing compression levels)
    quality_presets = {
        "low": "/printer",   # Low compression / High quality (300 DPI)
        "medium": "/ebook",  # Medium compression / Balanced quality (150 DPI)
        "high": "/screen"    # High compression / Maximum reduction (72/60 DPI)
    }
    preset = quality_presets.get(quality.lower(), "/ebook")
    
    # Setup output paths
    output_dir = os.path.dirname(file_path)
    output_filename = f"compressed_{uuid.uuid4()}.pdf"
    output_path = os.path.join(output_dir, output_filename)

    # 50% - Core Ghostscript run
    self.update_state(state="PROGRESS", meta={"progress": 50, "message": "Compressing PDF with Ghostscript..."})
    
    # Build GS arguments
    gs_cmd = [
        "gs",
        "-sDEVICE=pdfwrite",
        "-dCompatibilityLevel=1.4",
        f"-dPDFSETTINGS={preset}",
        "-dNOPAUSE",
        "-dQUIET",
        "-dBATCH",
    ]
    
    if quality.lower() == "medium":
        # Moderate downsampling and DCT JPEG compression for balanced ~60% size reduction
        gs_cmd.extend([
            "-dDownsampleColorImages=true",
            "-dColorImageDownsampleType=/Bicubic",
            "-dColorImageResolution=150",
            "-dDownsampleGrayImages=true",
            "-dGrayImageDownsampleType=/Bicubic",
            "-dGrayImageResolution=150",
            "-dDownsampleMonoImages=true",
            "-dMonoImageDownsampleType=/Bicubic",
            "-dMonoImageResolution=150",
            "-dColorImageFilter=/DCTEncode",
            "-dGrayImageFilter=/DCTEncode",
        ])
    elif quality.lower() == "high":
        # Aggressive downsampling and DCT JPEG compression for maximum ~80% size reduction
        gs_cmd.extend([
            "-dDownsampleColorImages=true",
            "-dColorImageDownsampleType=/Bicubic",
            "-dColorImageResolution=60",
            "-dDownsampleGrayImages=true",
            "-dGrayImageDownsampleType=/Bicubic",
            "-dGrayImageResolution=60",
            "-dDownsampleMonoImages=true",
            "-dMonoImageDownsampleType=/Bicubic",
            "-dMonoImageResolution=100",
            "-dColorImageFilter=/DCTEncode",
            "-dGrayImageFilter=/DCTEncode",
        ])
        
    gs_cmd.extend([
        f"-sOutputFile={output_path}",
        file_path
    ])
    
    try:
        # Get original file size
        original_size = os.path.getsize(file_path)
        
        logger.info(f"Running command: {' '.join(gs_cmd)}")
        result = subprocess.run(gs_cmd, capture_output=True, text=True, check=True)
        
        if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
            raise Exception("Ghostscript executed but output file was not created or is empty.")
            
        compressed_size = os.path.getsize(output_path)
        
        # 80% - Store processed file
        self.update_state(state="PROGRESS", meta={"progress": 80, "message": "Saving processed PDF file..."})
        
        original_name = None
        if original_filename:
            base, ext = os.path.splitext(original_filename)
            original_name = f"{base}_compressed{ext}"
        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store the compressed PDF.")
            
        # If stored locally (S3 not configured), schedule automatic deletion in 15 minutes (900 seconds)
        if not is_s3_configured():
            app_dir = os.path.dirname(os.path.abspath(__file__))
            backend_root = os.path.dirname(app_dir)
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)
            logger.info(f"Scheduled auto-cleanup for {local_dest_path} in 15 minutes.")
            
        logger.info("PDF compression completed successfully.")
        
        return {
            "status": "success",
            "original_size_bytes": original_size,
            "compressed_size_bytes": compressed_size,
            "download_url": download_url
        }
        
    except subprocess.CalledProcessError as e:
        logger.error(f"Ghostscript failed (exit code {e.returncode}): {e.stderr}")
        raise Exception(f"Ghostscript compression failed: {e.stderr}")
    except Exception as e:
        logger.error(f"Task failed: {str(e)}")
        raise e
    finally:
        # Clean up both the input and output local files
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                    logger.info(f"Cleaned up local file: {path}")
                except Exception as cleanup_err:
                    logger.warning(f"Failed to remove {path}: {cleanup_err}")

@celery.task
def merge_pdfs_task(file_paths: list[str], original_filename: str = None):
    """
    Combines a list of PDF files into a single output PDF, uploads it, and cleans up.
    """
    logger.info(f"Starting PDF merge task for files: {file_paths}")
    
    app_dir = os.path.dirname(os.path.abspath(__file__))
    backend_root = os.path.dirname(app_dir)
    
    # Resolve relative paths
    resolved_paths = []
    for fp in file_paths:
        fp = fp.replace('\\', '/')
        if not os.path.isabs(fp):
            fp = os.path.abspath(os.path.join(backend_root, fp))
        resolved_paths.append(fp)
        
    output_filename = f"merged_{uuid.uuid4()}.pdf"
    tmp_dir = os.path.join(backend_root, "tmp")
    os.makedirs(tmp_dir, exist_ok=True)
    output_path = os.path.join(tmp_dir, output_filename)
    
    merger = PdfWriter()
    try:
        # Append files in order
        for fp in resolved_paths:
            if os.path.exists(fp):
                merger.append(fp)
            else:
                logger.error(f"Merge file not found: {fp}")
                raise FileNotFoundError(f"File to merge not found: {fp}")
                
        with open(output_path, "wb") as f:
            merger.write(f)
        merger.close()
        
        # Save output using storage manager
        original_name = None
        if original_filename:
            base, ext = os.path.splitext(original_filename)
            original_name = f"{base}_merged{ext}"
        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store merged PDF.")
            
        # Schedule cleanup if local
        if not is_s3_configured():
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)
            logger.info(f"Scheduled merge cleanup for {local_dest_path}")
            
        return {
            "status": "success",
            "download_url": download_url
        }
        
    except Exception as e:
        logger.error(f"Merge task failed: {e}")
        raise e
    finally:
        # Cleanup merger if not closed
        try:
            merger.close()
        except Exception:
            pass
        # Cleanup input files
        for fp in resolved_paths:
            if fp and os.path.exists(fp):
                try:
                    os.remove(fp)
                except Exception as err:
                    logger.warning(f"Failed to remove input PDF: {err}")
        # Cleanup temporary output
        if os.path.exists(output_path):
            try:
                os.remove(output_path)
            except Exception:
                pass

@celery.task
def split_pdf_task(file_path: str, mode: str, pages_spec: str, original_filename: str = None):
    """
    Splits a PDF by custom page range (returning a PDF) or every page (returning a ZIP).
    """
    logger.info(f"Starting split task. Mode: {mode}, Spec: {pages_spec}")
    
    app_dir = os.path.dirname(os.path.abspath(__file__))
    backend_root = os.path.dirname(app_dir)
    
    file_path = file_path.replace('\\', '/')
    if not os.path.isabs(file_path):
        file_path = os.path.abspath(os.path.join(backend_root, file_path))
        
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File to split not found: {file_path}")
        
    tmp_dir = os.path.join(backend_root, "tmp")
    os.makedirs(tmp_dir, exist_ok=True)
    
    output_filename = ""
    output_path = ""
    
    try:
        reader = PdfReader(file_path)
        total_pages = len(reader.pages)
        
        if mode.lower() == "all" and total_pages > 1:
            # Split every page into separate PDFs inside a ZIP
            output_filename = f"split_{uuid.uuid4()}.zip"
            output_path = os.path.join(tmp_dir, output_filename)
            
            with zipfile.ZipFile(output_path, "w") as zipf:
                for page_num in range(total_pages):
                    writer = PdfWriter()
                    writer.add_page(reader.pages[page_num])
                    
                    page_name = f"page_{page_num + 1}.pdf"
                    page_temp_path = os.path.join(tmp_dir, page_name)
                    
                    with open(page_temp_path, "wb") as f:
                        writer.write(f)
                    
                    zipf.write(page_temp_path, page_name)
                    os.remove(page_temp_path)
                    
            original_name = None
            if original_filename:
                base, ext = os.path.splitext(original_filename)
                original_name = f"{base}_split.zip"
        else:
            # Custom range extraction into a single PDF, or single page PDF split
            if mode.lower() == "all":
                # Only 1 page exists
                pages_to_extract = [0]
            else:
                pages_to_extract = []
                for part in pages_spec.split(","):
                    part = part.strip()
                    if not part:
                        continue
                    if "-" in part:
                        start_str, end_str = part.split("-")
                        pages_to_extract.extend(range(int(start_str) - 1, int(end_str)))
                    else:
                        pages_to_extract.append(int(part) - 1)
            
            writer = PdfWriter()
            for page_num in pages_to_extract:
                if 0 <= page_num < total_pages:
                    writer.add_page(reader.pages[page_num])
                    
            output_filename = f"split_{uuid.uuid4()}.pdf"
            output_path = os.path.join(tmp_dir, output_filename)
            with open(output_path, "wb") as f:
                writer.write(f)
                
            original_name = None
            if original_filename:
                base, ext = os.path.splitext(original_filename)
                original_name = f"{base}_split.pdf"
        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store split PDF results.")
            
        # Schedule cleanup
        if not is_s3_configured():
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)
            
        return {
            "status": "success",
            "download_url": download_url
        }
        
    except Exception as e:
        logger.error(f"Split task failed: {e}")
        raise e
    finally:
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass

@celery.task
def pdf_to_images_task(file_path: str, img_format: str, dpi: int, original_filename: str = None):
    """
    Converts PDF pages into image formats (PNG/JPG) and packages them as a ZIP.
    """
    logger.info(f"Starting PDF to Image task. Format: {img_format}, DPI: {dpi}")
    
    app_dir = os.path.dirname(os.path.abspath(__file__))
    backend_root = os.path.dirname(app_dir)
    
    file_path = file_path.replace('\\', '/')
    if not os.path.isabs(file_path):
        file_path = os.path.abspath(os.path.join(backend_root, file_path))
        
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
        
    tmp_dir = os.path.join(backend_root, "tmp")
    os.makedirs(tmp_dir, exist_ok=True)
    
    output_filename = ""
    output_path = ""
    
    try:
        # Convert pages to PIL Images
        images = convert_from_path(file_path, dpi=dpi)
        
        if len(images) == 1:
            output_filename = f"image_{uuid.uuid4()}.{img_format.lower()}"
            output_path = os.path.join(tmp_dir, output_filename)
            images[0].save(output_path, format=img_format.upper())
            
            original_name = None
            if original_filename:
                base, ext = os.path.splitext(original_filename)
                original_name = f"{base}.{img_format.lower()}"
        else:
            output_filename = f"images_{uuid.uuid4()}.zip"
            output_path = os.path.join(tmp_dir, output_filename)
            with zipfile.ZipFile(output_path, "w") as zipf:
                for idx, img in enumerate(images):
                    img_name = f"page_{idx + 1}.{img_format.lower()}"
                    temp_img_path = os.path.join(tmp_dir, img_name)
                    
                    img.save(temp_img_path, format=img_format.upper())
                    zipf.write(temp_img_path, img_name)
                    os.remove(temp_img_path)
                    
            original_name = None
            if original_filename:
                base, ext = os.path.splitext(original_filename)
                original_name = f"{base}_images.zip"
                
        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store converted images ZIP.")
            
        # Schedule cleanup
        if not is_s3_configured():
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)
            
        return {
            "status": "success",
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"PDF to Image conversion failed: {e}")
        raise e
    finally:
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass

@celery.task
def images_to_pdf_task(file_paths: list[str], original_filename: str = None):
    """
    Compiles multiple images into a single PDF losslessly and cleans up.
    """
    logger.info(f"Starting Images to PDF task for files: {file_paths}")
    
    app_dir = os.path.dirname(os.path.abspath(__file__))
    backend_root = os.path.dirname(app_dir)
    
    # Resolve relative paths
    resolved_paths = []
    for fp in file_paths:
        fp = fp.replace('\\', '/')
        if not os.path.isabs(fp):
            fp = os.path.abspath(os.path.join(backend_root, fp))
        resolved_paths.append(fp)
        
    output_filename = f"compiled_{uuid.uuid4()}.pdf"
    tmp_dir = os.path.join(backend_root, "tmp")
    os.makedirs(tmp_dir, exist_ok=True)
    output_path = os.path.join(tmp_dir, output_filename)
    
    try:
        # Filter existing paths
        valid_paths = [fp for fp in resolved_paths if os.path.exists(fp)]
        if not valid_paths:
            raise Exception("No valid image files exist to compile.")
            
        # Convert losslessly with img2pdf
        pdf_bytes = img2pdf.convert(valid_paths)
        with open(output_path, "wb") as f:
            f.write(pdf_bytes)
            
        original_name = None
        if original_filename:
            base, ext = os.path.splitext(original_filename)
            original_name = f"{base}_compiled.pdf"
        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store compiled images PDF.")
            
        # Schedule cleanup
        if not is_s3_configured():
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)
            
        return {
            "status": "success",
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"Images to PDF task failed: {e}")
        raise e
    finally:
        # Cleanup input images
        for fp in resolved_paths:
            if fp and os.path.exists(fp):
                try:
                    os.remove(fp)
                except Exception:
                    pass
        # Cleanup temporary output
        if os.path.exists(output_path):
            try:
                os.remove(output_path)
            except Exception:
                pass

@celery.task
def extract_text_task(file_path: str, export_txt: bool, original_filename: str = None):
    """
    Extracts text from a PDF. Saves to a .txt file if export_txt is True, otherwise returns text snippet.
    """
    logger.info(f"Starting text extraction task. export_txt: {export_txt}")
    
    app_dir = os.path.dirname(os.path.abspath(__file__))
    backend_root = os.path.dirname(app_dir)
    
    file_path = file_path.replace('\\', '/')
    if not os.path.isabs(file_path):
        file_path = os.path.abspath(os.path.join(backend_root, file_path))
        
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
        
    output_filename = ""
    output_path = ""
    
    try:
        reader = PdfReader(file_path)
        extracted_content = ""
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                extracted_content += f"--- Page {i + 1} ---\n{text}\n\n"
                
        if not extracted_content:
            extracted_content = "No text could be extracted from this document."
            
        if export_txt:
            output_filename = f"extracted_{uuid.uuid4()}.txt"
            tmp_dir = os.path.join(backend_root, "tmp")
            os.makedirs(tmp_dir, exist_ok=True)
            output_path = os.path.join(tmp_dir, output_filename)
            
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(extracted_content)
                
            original_name = None
            if original_filename:
                base, ext = os.path.splitext(original_filename)
                original_name = f"{base}_extracted.txt"
            download_url = store_processed_file(output_path, output_filename, original_name=original_name)
            if not download_url:
                raise Exception("Failed to store extracted text file.")
                
            # Schedule cleanup
            if not is_s3_configured():
                local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
                delete_local_file_task.apply_async((local_dest_path,), countdown=900)
                
            return {
                "status": "success",
                "download_url": download_url,
                "text_snippet": extracted_content[:1500]
            }
        else:
            return {
                "status": "success",
                "extracted_text": extracted_content
            }
    except Exception as e:
        logger.error(f"Text extraction failed: {e}")
        raise e
    finally:
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass


@celery.task(bind=True)
def pdf_to_docx_task(self, file_path: str, original_filename: str = None):
    """
    Converts a PDF file into a DOCX file using pdf2docx, uploads it, and cleans up.
    """
    logger.info(f"Starting PDF to DOCX task for {file_path}")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Initializing conversion..."})
    
    file_path = file_path.replace('\\', '/')
    if not os.path.isabs(file_path):
        app_dir = os.path.dirname(os.path.abspath(__file__))
        backend_root = os.path.dirname(app_dir)
        file_path = os.path.abspath(os.path.join(backend_root, file_path))
        
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")
        
    output_dir = os.path.dirname(file_path)
    output_filename = f"converted_{uuid.uuid4()}.docx"
    output_path = os.path.join(output_dir, output_filename)
    
    try:
        self.update_state(state="PROGRESS", meta={"progress": 50, "message": "Parsing layout and converting to Word..."})
        
        cv = Converter(file_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()
        
        self.update_state(state="PROGRESS", meta={"progress": 85, "message": "Saving converted file..."})
        
        original_name = None
        if original_filename:
            base, ext = os.path.splitext(original_filename)
            original_name = f"{base}.docx"
            
        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store converted DOCX file.")
            
        if not is_s3_configured():
            app_dir = os.path.dirname(os.path.abspath(__file__))
            backend_root = os.path.dirname(app_dir)
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)
            
        return {
            "status": "success",
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"PDF to DOCX conversion failed: {e}")
        raise e
    finally:
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass


@celery.task(bind=True)
def docx_to_pdf_task(self, file_path: str, original_filename: str = None):
    """
    Converts a DOCX file into a PDF file using LibreOffice (or MS Word on Windows), uploads it, and cleans up.
    """
    logger.info(f"Starting DOCX to PDF task for {file_path}")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Initializing conversion..."})
    
    file_path = file_path.replace('\\', '/')
    if not os.path.isabs(file_path):
        app_dir = os.path.dirname(os.path.abspath(__file__))
        backend_root = os.path.dirname(app_dir)
        file_path = os.path.abspath(os.path.join(backend_root, file_path))
        
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")
        
    output_dir = os.path.dirname(file_path)
    output_filename = f"converted_{uuid.uuid4()}.pdf"
    output_path = os.path.join(output_dir, output_filename)
    
    try:
        self.update_state(state="PROGRESS", meta={"progress": 50, "message": "Converting document to PDF..."})
        
        # Determine conversion command / method
        soffice_cmd = None
        if sys.platform == "win32":
            # Common Windows paths for LibreOffice
            paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
            ]
            for p in paths:
                if os.path.exists(p):
                    soffice_cmd = p
                    break
                    
        if not soffice_cmd:
            # Check PATH
            for cmd in ["soffice", "libreoffice"]:
                if shutil.which(cmd):
                    soffice_cmd = cmd
                    break
                    
        converted = False
        if soffice_cmd:
            # Run LibreOffice headless conversion
            cmd = [soffice_cmd, "--headless", "--convert-to", "pdf", "--outdir", output_dir, file_path]
            logger.info(f"Running LibreOffice conversion: {' '.join(cmd)}")
            result = subprocess.run(cmd, capture_output=True, text=True, check=True)
            
            # LibreOffice saves output in output_dir as [original_base_name].pdf
            base_name = os.path.splitext(os.path.basename(file_path))[0]
            generated_pdf = os.path.join(output_dir, f"{base_name}.pdf")
            if os.path.exists(generated_pdf):
                if generated_pdf != output_path:
                    shutil.move(generated_pdf, output_path)
                converted = True
                
        if not converted and sys.platform == "win32":
            # Fall back to MS Word COM Automation on Windows
            try:
                import win32com.client
                logger.info("Attempting MS Word COM fallback on Windows")
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                doc = word.Documents.Open(os.path.abspath(file_path))
                # 17 is wdFormatPDF
                doc.SaveAs(os.path.abspath(output_path), FileFormat=17)
                doc.Close()
                word.Quit()
                converted = True
            except Exception as com_err:
                logger.error(f"MS Word COM fallback failed: {com_err}")
                
        if not converted:
            raise Exception("Could not convert DOCX to PDF. Neither LibreOffice headless nor MS Word is available.")
            
        self.update_state(state="PROGRESS", meta={"progress": 85, "message": "Saving converted PDF..."})
        
        original_name = None
        if original_filename:
            base, ext = os.path.splitext(original_filename)
            original_name = f"{base}.pdf"
            
        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store converted PDF file.")
            
        if not is_s3_configured():
            app_dir = os.path.dirname(os.path.abspath(__file__))
            backend_root = os.path.dirname(app_dir)
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)
            
        return {
            "status": "success",
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"DOCX to PDF conversion failed: {e}")
        raise e
    finally:
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass


def create_watermark_pdf(output_path, text, color="gray", opacity=0.3, rotation=45):
    from reportlab.pdfgen import canvas
    from reportlab.lib.colors import HexColor
    from reportlab.lib.pagesizes import letter

    c = canvas.Canvas(output_path, pagesize=letter)
    width, height = letter
    c.translate(width / 2.0, height / 2.0)
    c.rotate(rotation)
    c.setFillAlpha(opacity)
    
    colors_map = {
        "red": "#FF0000",
        "green": "#00FF00",
        "blue": "#0000FF",
        "gray": "#808080",
        "black": "#000000"
    }
    hex_color = colors_map.get(color.lower(), "#808080")
    c.setFillColor(HexColor(hex_color))
    c.setFont("Helvetica", 48)
    c.drawCentredString(0, 0, text)
    c.save()


def create_page_number_pdf(output_path, width, height, current_page, total_pages, position):
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(output_path, pagesize=(width, height))
    c.setFont("Helvetica", 10)
    c.setFillColorRGB(0.5, 0.5, 0.5)

    text = f"Page {current_page} of {total_pages}"
    margin = 36
    x = width / 2.0
    y = margin

    if "top" in position:
        y = height - margin

    if "left" in position:
        x = margin
        align = "left"
    elif "right" in position:
        x = width - margin
        align = "right"
    else:
        x = width / 2.0
        align = "center"

    if align == "left":
        c.drawString(x, y, text)
    elif align == "right":
        c.drawRightString(x, y, text)
    else:
        c.drawCentredString(x, y, text)

    c.save()


@celery.task(bind=True)
def protect_pdf_task(self, file_path: str, password: str, original_filename: str = None):
    """
    Encrypts a PDF with a password.
    """
    logger.info(f"Starting PDF encryption task for {file_path}")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Initializing encryption..."})

    file_path = file_path.replace('\\', '/')
    if not os.path.isabs(file_path):
        app_dir = os.path.dirname(os.path.abspath(__file__))
        backend_root = os.path.dirname(app_dir)
        file_path = os.path.abspath(os.path.join(backend_root, file_path))

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")

    output_dir = os.path.dirname(file_path)
    output_filename = f"protected_{uuid.uuid4()}.pdf"
    output_path = os.path.join(output_dir, output_filename)

    try:
        self.update_state(state="PROGRESS", meta={"progress": 50, "message": "Encrypting PDF document..."})
        reader = PdfReader(file_path)
        writer = PdfWriter()
        writer.append(reader)
        writer.encrypt(password)

        with open(output_path, "wb") as f:
            writer.write(f)

        self.update_state(state="PROGRESS", meta={"progress": 85, "message": "Saving encrypted PDF..."})
        original_name = None
        if original_filename:
            base, ext = os.path.splitext(original_filename)
            original_name = f"{base}_protected{ext}"

        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store encrypted PDF.")

        if not is_s3_configured():
            app_dir = os.path.dirname(os.path.abspath(__file__))
            backend_root = os.path.dirname(app_dir)
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)

        return {
            "status": "success",
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"Encryption task failed: {e}")
        raise e
    finally:
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass


@celery.task(bind=True)
def unlock_pdf_task(self, file_path: str, password: str, original_filename: str = None):
    """
    Decrypts/removes password protection from a PDF.
    """
    logger.info(f"Starting PDF decryption task for {file_path}")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Initializing decryption..."})

    file_path = file_path.replace('\\', '/')
    if not os.path.isabs(file_path):
        app_dir = os.path.dirname(os.path.abspath(__file__))
        backend_root = os.path.dirname(app_dir)
        file_path = os.path.abspath(os.path.join(backend_root, file_path))

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")

    output_dir = os.path.dirname(file_path)
    output_filename = f"unlocked_{uuid.uuid4()}.pdf"
    output_path = os.path.join(output_dir, output_filename)

    try:
        self.update_state(state="PROGRESS", meta={"progress": 50, "message": "Decrypting PDF document..."})
        reader = PdfReader(file_path)
        if reader.is_encrypted:
            # Try decrypting
            result = reader.decrypt(password)
            if result == 0:
                raise Exception("Incorrect password provided for the PDF file.")

        writer = PdfWriter()
        writer.append(reader)

        with open(output_path, "wb") as f:
            writer.write(f)

        self.update_state(state="PROGRESS", meta={"progress": 85, "message": "Saving unlocked PDF..."})
        original_name = None
        if original_filename:
            base, ext = os.path.splitext(original_filename)
            original_name = f"{base}_unlocked{ext}"

        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store unlocked PDF.")

        if not is_s3_configured():
            app_dir = os.path.dirname(os.path.abspath(__file__))
            backend_root = os.path.dirname(app_dir)
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)

        return {
            "status": "success",
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"Decryption task failed: {e}")
        raise e
    finally:
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass


@celery.task(bind=True)
def rotate_pdf_task(self, file_path: str, rotation: int, original_filename: str = None):
    """
    Rotates all pages in a PDF.
    """
    logger.info(f"Starting PDF rotation task for {file_path}. Rotation angle: {rotation}")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Initializing rotation..."})

    file_path = file_path.replace('\\', '/')
    if not os.path.isabs(file_path):
        app_dir = os.path.dirname(os.path.abspath(__file__))
        backend_root = os.path.dirname(app_dir)
        file_path = os.path.abspath(os.path.join(backend_root, file_path))

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")

    output_dir = os.path.dirname(file_path)
    output_filename = f"rotated_{uuid.uuid4()}.pdf"
    output_path = os.path.join(output_dir, output_filename)

    try:
        self.update_state(state="PROGRESS", meta={"progress": 50, "message": "Rotating PDF pages..."})
        reader = PdfReader(file_path)
        writer = PdfWriter()

        for page in reader.pages:
            page.rotate(rotation)
            writer.add_page(page)

        with open(output_path, "wb") as f:
            writer.write(f)

        self.update_state(state="PROGRESS", meta={"progress": 85, "message": "Saving rotated PDF..."})
        original_name = None
        if original_filename:
            base, ext = os.path.splitext(original_filename)
            original_name = f"{base}_rotated{ext}"

        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store rotated PDF.")

        if not is_s3_configured():
            app_dir = os.path.dirname(os.path.abspath(__file__))
            backend_root = os.path.dirname(app_dir)
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)

        return {
            "status": "success",
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"Rotation task failed: {e}")
        raise e
    finally:
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass


@celery.task(bind=True)
def organize_pdf_task(self, file_path: str, page_order: str, original_filename: str = None):
    """
    Reorders or deletes specific pages in a PDF.
    """
    logger.info(f"Starting PDF organization task for {file_path}. Order: {page_order}")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Initializing organization..."})

    file_path = file_path.replace('\\', '/')
    if not os.path.isabs(file_path):
        app_dir = os.path.dirname(os.path.abspath(__file__))
        backend_root = os.path.dirname(app_dir)
        file_path = os.path.abspath(os.path.join(backend_root, file_path))

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")

    output_dir = os.path.dirname(file_path)
    output_filename = f"organized_{uuid.uuid4()}.pdf"
    output_path = os.path.join(output_dir, output_filename)

    try:
        self.update_state(state="PROGRESS", meta={"progress": 50, "message": "Rebuilding PDF structure..."})
        reader = PdfReader(file_path)
        writer = PdfWriter()

        indices = [int(p.strip()) - 1 for p in page_order.split(",") if p.strip()]
        total_pages = len(reader.pages)

        for idx in indices:
            if 0 <= idx < total_pages:
                writer.add_page(reader.pages[idx])

        with open(output_path, "wb") as f:
            writer.write(f)

        self.update_state(state="PROGRESS", meta={"progress": 85, "message": "Saving reorganized PDF..."})
        original_name = None
        if original_filename:
            base, ext = os.path.splitext(original_filename)
            original_name = f"{base}_organized{ext}"

        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store organized PDF.")

        if not is_s3_configured():
            app_dir = os.path.dirname(os.path.abspath(__file__))
            backend_root = os.path.dirname(app_dir)
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)

        return {
            "status": "success",
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"Organization task failed: {e}")
        raise e
    finally:
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass


@celery.task(bind=True)
def watermark_pdf_task(self, file_path: str, text: str, color: str = "gray", opacity: float = 0.3, rotation: int = 45, original_filename: str = None):
    """
    Applies a text watermark onto all pages of a PDF.
    """
    logger.info(f"Starting PDF watermarking task for {file_path}")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Initializing watermark..."})

    file_path = file_path.replace('\\', '/')
    if not os.path.isabs(file_path):
        app_dir = os.path.dirname(os.path.abspath(__file__))
        backend_root = os.path.dirname(app_dir)
        file_path = os.path.abspath(os.path.join(backend_root, file_path))

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")

    output_dir = os.path.dirname(file_path)
    output_filename = f"watermarked_{uuid.uuid4()}.pdf"
    output_path = os.path.join(output_dir, output_filename)

    try:
        self.update_state(state="PROGRESS", meta={"progress": 50, "message": "Drawing watermark template..."})
        
        # Create temp watermark PDF
        temp_watermark_pdf = os.path.join(output_dir, f"temp_wm_{uuid.uuid4()}.pdf")
        create_watermark_pdf(temp_watermark_pdf, text, color, opacity, rotation)

        self.update_state(state="PROGRESS", meta={"progress": 70, "message": "Merging watermark into document..."})
        reader = PdfReader(file_path)
        wm_reader = PdfReader(temp_watermark_pdf)
        wm_page = wm_reader.pages[0]

        writer = PdfWriter()
        for page in reader.pages:
            # We scale or merge the watermark page over the original page
            page.merge_page(wm_page)
            writer.add_page(page)

        with open(output_path, "wb") as f:
            writer.write(f)

        # Clean up temp watermark file
        if os.path.exists(temp_watermark_pdf):
            os.remove(temp_watermark_pdf)

        self.update_state(state="PROGRESS", meta={"progress": 85, "message": "Saving watermarked PDF..."})
        original_name = None
        if original_filename:
            base, ext = os.path.splitext(original_filename)
            original_name = f"{base}_watermarked{ext}"

        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store watermarked PDF.")

        if not is_s3_configured():
            app_dir = os.path.dirname(os.path.abspath(__file__))
            backend_root = os.path.dirname(app_dir)
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)

        return {
            "status": "success",
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"Watermark task failed: {e}")
        raise e
    finally:
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass


@celery.task(bind=True)
def add_page_numbers_task(self, file_path: str, position: str = "bottom-center", original_filename: str = None):
    """
    Stamps page numbers onto all pages of a PDF.
    """
    logger.info(f"Starting page numbering task for {file_path}. Position: {position}")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Initializing page numbers..."})

    file_path = file_path.replace('\\', '/')
    if not os.path.isabs(file_path):
        app_dir = os.path.dirname(os.path.abspath(__file__))
        backend_root = os.path.dirname(app_dir)
        file_path = os.path.abspath(os.path.join(backend_root, file_path))

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")

    output_dir = os.path.dirname(file_path)
    output_filename = f"numbered_{uuid.uuid4()}.pdf"
    output_path = os.path.join(output_dir, output_filename)

    try:
        self.update_state(state="PROGRESS", meta={"progress": 50, "message": "Numbering PDF pages..."})
        reader = PdfReader(file_path)
        writer = PdfWriter()
        total_pages = len(reader.pages)

        for idx, page in enumerate(reader.pages):
            box = page.mediabox
            width = float(box.width)
            height = float(box.height)

            # Generate dynamic single-page canvas with page number text
            temp_pagenum_pdf = os.path.join(output_dir, f"temp_num_{uuid.uuid4()}.pdf")
            create_page_number_pdf(temp_pagenum_pdf, width, height, idx + 1, total_pages, position)

            # Merge and write
            num_reader = PdfReader(temp_pagenum_pdf)
            page.merge_page(num_reader.pages[0])
            writer.add_page(page)

            # Clean up temp page PDF
            if os.path.exists(temp_pagenum_pdf):
                os.remove(temp_pagenum_pdf)

        with open(output_path, "wb") as f:
            writer.write(f)

        self.update_state(state="PROGRESS", meta={"progress": 85, "message": "Saving numbered PDF..."})
        original_name = None
        if original_filename:
            base, ext = os.path.splitext(original_filename)
            original_name = f"{base}_numbered{ext}"

        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store numbered PDF.")

        if not is_s3_configured():
            app_dir = os.path.dirname(os.path.abspath(__file__))
            backend_root = os.path.dirname(app_dir)
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)

        return {
            "status": "success",
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"Page numbering task failed: {e}")
        raise e
    finally:
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass


@celery.task(bind=True)
def ocr_pdf_task(self, file_path: str, language: str = "eng", original_filename: str = None):
    """
    Performs OCR on a PDF document using ocrmypdf CLI, uploads it, and cleans up.
    """
    logger.info(f"Starting OCR task for {file_path}. Language: {language}")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Initializing OCR engine..."})

    file_path = file_path.replace('\\', '/')
    if not os.path.isabs(file_path):
        app_dir = os.path.dirname(os.path.abspath(__file__))
        backend_root = os.path.dirname(app_dir)
        file_path = os.path.abspath(os.path.join(backend_root, file_path))

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")

    output_dir = os.path.dirname(file_path)
    output_filename = f"ocr_{uuid.uuid4()}.pdf"
    output_path = os.path.join(output_dir, output_filename)

    try:
        self.update_state(state="PROGRESS", meta={"progress": 50, "message": "Performing OCR with Tesseract... (this may take a while)"})
        
        # Run ocrmypdf command
        # --skip-text skips OCR on pages that already have text (highly recommended for speed/correctness)
        cmd = ["ocrmypdf", "--skip-text", "-l", language, file_path, output_path]
        logger.info(f"Running OCR command: {' '.join(cmd)}")
        result = subprocess.run(cmd, capture_output=True, text=True, check=True)

        self.update_state(state="PROGRESS", meta={"progress": 85, "message": "Saving OCR-processed PDF..."})
        original_name = None
        if original_filename:
            base, ext = os.path.splitext(original_filename)
            original_name = f"{base}_ocr{ext}"

        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store OCR-processed PDF.")

        if not is_s3_configured():
            app_dir = os.path.dirname(os.path.abspath(__file__))
            backend_root = os.path.dirname(app_dir)
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)

        return {
            "status": "success",
            "download_url": download_url
        }
    except FileNotFoundError:
        logger.error("ocrmypdf command not found in PATH.")
        raise Exception("OCR tool is not installed or configured on the server. Requires ocrmypdf and Tesseract OCR.")
    except subprocess.CalledProcessError as e:
        logger.error(f"OCR command failed (code {e.returncode}): {e.stderr}")
        raise Exception(f"OCR processing failed: {e.stderr or e.stdout}")
    except Exception as e:
        logger.error(f"OCR task failed: {e}")
        raise e
    finally:
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass


def _resolve_path(path_str: str) -> str:
    path_str = path_str.replace('\\', '/')
    if not os.path.isabs(path_str):
        app_dir = os.path.dirname(os.path.abspath(__file__))
        backend_root = os.path.dirname(app_dir)
        path_str = os.path.abspath(os.path.join(backend_root, path_str))
    return path_str


@celery.task(bind=True)
def compress_image_task(self, file_path: str, quality: int, original_filename: str = None):
    """
    Compresses an image (JPEG, PNG, WEBP) using Pillow.
    """
    logger.info(f"Starting image compression task for {file_path} with quality={quality}")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Initializing image compression..."})
    
    file_path = _resolve_path(file_path)
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")
        
    ext = os.path.splitext(file_path)[1].lower()
    if ext not in [".jpg", ".jpeg", ".png", ".webp"]:
        raise ValueError(f"Unsupported image format: {ext}")
        
    output_filename = f"compressed_{uuid.uuid4()}{ext}"
    output_dir = os.path.dirname(file_path)
    output_path = os.path.join(output_dir, output_filename)
    
    try:
        self.update_state(state="PROGRESS", meta={"progress": 50, "message": "Applying compression..."})
        with Image.open(file_path) as img:
            # JPEG does not support alpha channel (RGBA), convert to RGB
            if ext in [".jpg", ".jpeg"] and img.mode in ("RGBA", "LA"):
                img = img.convert("RGB")
                
            if ext in [".jpg", ".jpeg"]:
                img.save(output_path, "JPEG", quality=quality, optimize=True)
            elif ext == ".webp":
                img.save(output_path, "WEBP", quality=quality, optimize=True)
            elif ext == ".png":
                # For PNG, quality doesn't apply directly. We use optimize=True.
                # If quality is low/medium, we can quantize colors to reduce size (reduce to 8-bit palette).
                if quality < 50:
                    img = img.quantize(colors=128).convert("RGBA")
                elif quality < 80:
                    img = img.quantize(colors=256).convert("RGBA")
                img.save(output_path, "PNG", optimize=True)
                
        self.update_state(state="PROGRESS", meta={"progress": 80, "message": "Saving compressed image..."})
        
        original_name = None
        if original_filename:
            base, ext_orig = os.path.splitext(original_filename)
            original_name = f"{base}_compressed{ext_orig}"
            
        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store compressed image.")
            
        if not is_s3_configured():
            app_dir = os.path.dirname(os.path.abspath(__file__))
            backend_root = os.path.dirname(app_dir)
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)
            
        return {
            "status": "success",
            "original_size_bytes": os.path.getsize(file_path),
            "compressed_size_bytes": os.path.getsize(output_path),
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"Image compression failed: {e}")
        raise e
    finally:
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass


@celery.task(bind=True)
def resize_image_task(self, file_path: str, width: int = None, height: int = None, percentage: int = None, maintain_aspect: bool = True, original_filename: str = None):
    """
    Resizes an image by width/height pixels or by a percentage.
    """
    logger.info(f"Starting image resize task for {file_path}")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Initializing image resize..."})
    
    file_path = _resolve_path(file_path)
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")
        
    ext = os.path.splitext(file_path)[1].lower()
    output_filename = f"resized_{uuid.uuid4()}{ext}"
    output_dir = os.path.dirname(file_path)
    output_path = os.path.join(output_dir, output_filename)
    
    try:
        self.update_state(state="PROGRESS", meta={"progress": 50, "message": "Resizing image..."})
        with Image.open(file_path) as img:
            orig_w, orig_h = img.size
            
            # Determine new size
            if percentage is not None:
                new_w = int(orig_w * (percentage / 100.0))
                new_h = int(orig_h * (percentage / 100.0))
            else:
                if width and height:
                    if maintain_aspect:
                        # Scale to fit within width/height while maintaining aspect ratio
                        ratio = min(width / orig_w, height / orig_h)
                        new_w = int(orig_w * ratio)
                        new_h = int(orig_h * ratio)
                    else:
                        new_w = width
                        new_h = height
                elif width:
                    new_w = width
                    new_h = int(orig_h * (width / orig_w))
                elif height:
                    new_h = height
                    new_w = int(orig_w * (height / orig_h))
                else:
                    new_w, new_h = orig_w, orig_h
                    
            # Ensure dimensions are at least 1px
            new_w = max(1, new_w)
            new_h = max(1, new_h)
            
            resized_img = img.resize((new_w, new_h), Image.Resampling.LANCZOS)
            resized_img.save(output_path, save_all=True if ext == ".gif" else False)
            
        self.update_state(state="PROGRESS", meta={"progress": 80, "message": "Saving resized image..."})
        
        original_name = None
        if original_filename:
            base, ext_orig = os.path.splitext(original_filename)
            original_name = f"{base}_resized{ext_orig}"
            
        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store resized image.")
            
        if not is_s3_configured():
            app_dir = os.path.dirname(os.path.abspath(__file__))
            backend_root = os.path.dirname(app_dir)
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)
            
        return {
            "status": "success",
            "download_url": download_url,
            "width": new_w,
            "height": new_h
        }
    except Exception as e:
        logger.error(f"Image resize failed: {e}")
        raise e
    finally:
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass


@celery.task(bind=True)
def crop_image_task(self, file_path: str, x: int, y: int, width: int, height: int, original_filename: str = None):
    """
    Crops an image using bounding box coordinates (x, y, width, height).
    """
    logger.info(f"Starting image crop task for {file_path} at x={x}, y={y}, w={width}, h={height}")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Initializing image crop..."})
    
    file_path = _resolve_path(file_path)
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")
        
    ext = os.path.splitext(file_path)[1].lower()
    output_filename = f"cropped_{uuid.uuid4()}{ext}"
    output_dir = os.path.dirname(file_path)
    output_path = os.path.join(output_dir, output_filename)
    
    try:
        self.update_state(state="PROGRESS", meta={"progress": 50, "message": "Cropping image..."})
        with Image.open(file_path) as img:
            orig_w, orig_h = img.size
            
            # Ensure coordinates are within image bounds
            x = max(0, min(x, orig_w))
            y = max(0, min(y, orig_h))
            width = max(1, min(width, orig_w - x))
            height = max(1, min(height, orig_h - y))
            
            box = (x, y, x + width, y + height)
            cropped_img = img.crop(box)
            cropped_img.save(output_path)
            
        self.update_state(state="PROGRESS", meta={"progress": 80, "message": "Saving cropped image..."})
        
        original_name = None
        if original_filename:
            base, ext_orig = os.path.splitext(original_filename)
            original_name = f"{base}_cropped{ext_orig}"
            
        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store cropped image.")
            
        if not is_s3_configured():
            app_dir = os.path.dirname(os.path.abspath(__file__))
            backend_root = os.path.dirname(app_dir)
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)
            
        return {
            "status": "success",
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"Image crop failed: {e}")
        raise e
    finally:
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass


@celery.task(bind=True)
def convert_image_task(self, file_path: str, target_format: str, original_filename: str = None):
    """
    Converts image to target format (JPG, PNG, WEBP, BMP, TIFF).
    """
    logger.info(f"Starting image conversion task for {file_path} to format {target_format}")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Initializing image conversion..."})
    
    file_path = _resolve_path(file_path)
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")
        
    target_format = target_format.lower()
    valid_formats = {"jpg": "jpeg", "jpeg": "jpeg", "png": "png", "webp": "webp", "bmp": "bmp", "tiff": "tiff"}
    
    if target_format not in valid_formats:
        raise ValueError(f"Unsupported target format: {target_format}")
        
    pillow_format = valid_formats[target_format]
    out_ext = f".{target_format}"
    if out_ext == ".jpeg":
        out_ext = ".jpg"
        
    output_filename = f"converted_{uuid.uuid4()}{out_ext}"
    output_dir = os.path.dirname(file_path)
    output_path = os.path.join(output_dir, output_filename)
    
    try:
        self.update_state(state="PROGRESS", meta={"progress": 50, "message": "Converting format..."})
        with Image.open(file_path) as img:
            # Convert RGBA to RGB for JPEG
            if pillow_format == "jpeg" and img.mode in ("RGBA", "LA"):
                img = img.convert("RGB")
            elif img.mode == "P" and pillow_format not in ("png", "gif"):
                # Convert palette mode to RGB/RGBA for webp/tiff/bmp
                img = img.convert("RGBA" if "transparency" in img.info else "RGB")
                
            img.save(output_path, pillow_format.upper())
            
        self.update_state(state="PROGRESS", meta={"progress": 80, "message": "Saving converted image..."})
        
        original_name = None
        if original_filename:
            base, _ = os.path.splitext(original_filename)
            original_name = f"{base}{out_ext}"
            
        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store converted image.")
            
        if not is_s3_configured():
            app_dir = os.path.dirname(os.path.abspath(__file__))
            backend_root = os.path.dirname(app_dir)
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)
            
        return {
            "status": "success",
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"Image conversion failed: {e}")
        raise e
    finally:
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass


@celery.task(bind=True)
def rotate_image_task(self, file_path: str, angle: int, original_filename: str = None):
    """
    Rotates an image by 90, 180, or 270 degrees.
    """
    logger.info(f"Starting image rotation task for {file_path} with angle={angle}")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Initializing image rotation..."})
    
    file_path = _resolve_path(file_path)
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")
        
    ext = os.path.splitext(file_path)[1].lower()
    output_filename = f"rotated_{uuid.uuid4()}{ext}"
    output_dir = os.path.dirname(file_path)
    output_path = os.path.join(output_dir, output_filename)
    
    try:
        self.update_state(state="PROGRESS", meta={"progress": 50, "message": "Rotating image..."})
        with Image.open(file_path) as img:
            # Rotate using transpose for clean 90-degree step rotations
            if angle == 90:
                rotated_img = img.transpose(Image.Transpose.ROTATE_270)
            elif angle == 180:
                rotated_img = img.transpose(Image.Transpose.ROTATE_180)
            elif angle == 270:
                rotated_img = img.transpose(Image.Transpose.ROTATE_90)
            else:
                rotated_img = img.rotate(-angle, expand=True)
                
            rotated_img.save(output_path)
            
        self.update_state(state="PROGRESS", meta={"progress": 80, "message": "Saving rotated image..."})
        
        original_name = None
        if original_filename:
            base, ext_orig = os.path.splitext(original_filename)
            original_name = f"{base}_rotated{ext_orig}"
            
        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store rotated image.")
            
        if not is_s3_configured():
            app_dir = os.path.dirname(os.path.abspath(__file__))
            backend_root = os.path.dirname(app_dir)
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)
            
        return {
            "status": "success",
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"Image rotation failed: {e}")
        raise e
    finally:
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass


@celery.task(bind=True)
def watermark_image_task(self, file_path: str, text: str, color: str = "gray", opacity: float = 0.3, position: str = "center", original_filename: str = None):
    """
    Overlays a text watermark onto an image.
    Position can be: "center", "top-left", "top-right", "bottom-left", "bottom-right", "tile".
    """
    logger.info(f"Starting image watermark task for {file_path} with text='{text}'")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Initializing image watermark..."})
    
    file_path = _resolve_path(file_path)
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")
        
    ext = os.path.splitext(file_path)[1].lower()
    output_filename = f"watermarked_{uuid.uuid4()}{ext}"
    output_dir = os.path.dirname(file_path)
    output_path = os.path.join(output_dir, output_filename)
    
    try:
        self.update_state(state="PROGRESS", meta={"progress": 50, "message": "Applying watermark..."})
        with Image.open(file_path) as img:
            # We need RGBA mode for transparency/opacity layer
            watermark_layer = Image.new("RGBA", img.size, (0, 0, 0, 0))
            draw = ImageDraw.Draw(watermark_layer)
            
            colors_map = {
                "red": (255, 0, 0),
                "green": (0, 255, 0),
                "blue": (0, 0, 255),
                "gray": (128, 128, 128),
                "black": (0, 0, 0),
                "white": (255, 255, 255)
            }
            base_color = colors_map.get(color.lower(), (128, 128, 128))
            rgba_color = base_color + (int(opacity * 255),)
            
            # Choose a dynamic font size based on image size
            font_size = max(16, int(img.size[0] / 20))
            
            try:
                # Load default truetype font or fallback to default
                font = ImageFont.load_default()
                if sys.platform == "win32":
                    font = ImageFont.truetype("arial.ttf", font_size)
                else:
                    font = ImageFont.truetype("LiberationSans-Regular.ttf", font_size)
            except IOError:
                font = ImageFont.load_default()
                
            w, h = img.size
            
            if hasattr(draw, "textbbox"):
                bbox = draw.textbbox((0, 0), text, font=font)
                text_w = bbox[2] - bbox[0]
                text_h = bbox[3] - bbox[1]
            else:
                text_w, text_h = draw.textsize(text, font=font) if hasattr(draw, "textsize") else (font_size * len(text) * 0.6, font_size)
                
            # Place watermark based on position
            if position == "center":
                draw.text(((w - text_w) / 2, (h - text_h) / 2), text, font=font, fill=rgba_color)
            elif position == "top-left":
                draw.text((20, 20), text, font=font, fill=rgba_color)
            elif position == "top-right":
                draw.text((w - text_w - 20, 20), text, font=font, fill=rgba_color)
            elif position == "bottom-left":
                draw.text((20, h - text_h - 20), text, font=font, fill=rgba_color)
            elif position == "bottom-right":
                draw.text((w - text_w - 20, h - text_h - 20), text, font=font, fill=rgba_color)
            elif position == "tile":
                step_x = max(100, int(text_w * 1.5))
                step_y = max(100, int(text_h * 2.0))
                for x_pos in range(0, w, step_x):
                    for y_pos in range(0, h, step_y):
                        draw.text((x_pos, y_pos), text, font=font, fill=rgba_color)
                        
            orig_rgba = img.convert("RGBA")
            watermarked_img = Image.alpha_composite(orig_rgba, watermark_layer)
            
            if ext in [".jpg", ".jpeg"]:
                watermarked_img.convert("RGB").save(output_path, "JPEG")
            else:
                watermarked_img.save(output_path)
                
        self.update_state(state="PROGRESS", meta={"progress": 80, "message": "Saving watermarked image..."})
        
        original_name = None
        if original_filename:
            base, ext_orig = os.path.splitext(original_filename)
            original_name = f"{base}_watermark{ext_orig}"
            
        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store watermarked image.")
            
        if not is_s3_configured():
            app_dir = os.path.dirname(os.path.abspath(__file__))
            backend_root = os.path.dirname(app_dir)
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)
            
        return {
            "status": "success",
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"Image watermark failed: {e}")
        raise e
    finally:
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass


def _convert_docx_or_pptx_to_pdf(file_path: str, output_dir: str) -> str:
    """
    Converts a document (DOCX/PPTX) to PDF using LibreOffice headless mode.
    Returns the path to the generated PDF.
    """
    soffice_cmd = None
    if sys.platform == "win32":
        paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
        ]
        for p in paths:
            if os.path.exists(p):
                soffice_cmd = p
                break
                
    if not soffice_cmd:
        for cmd in ["soffice", "libreoffice"]:
            if shutil.which(cmd):
                soffice_cmd = cmd
                break
                
    if not soffice_cmd:
        raise Exception("LibreOffice soffice command not found. Cannot perform document conversion.")
        
    cmd = [soffice_cmd, "--headless", "--convert-to", "pdf", "--outdir", output_dir, file_path]
    logger.info(f"Running LibreOffice conversion command: {' '.join(cmd)}")
    subprocess.run(cmd, capture_output=True, text=True, check=True)
    
    base_name = os.path.splitext(os.path.basename(file_path))[0]
    generated_pdf = os.path.join(output_dir, f"{base_name}.pdf")
    if not os.path.exists(generated_pdf):
        raise FileNotFoundError(f"LibreOffice execution completed, but output PDF was not found at {generated_pdf}")
        
    return generated_pdf


@celery.task(bind=True)
def merge_docx_task(self, file_paths: list[str], original_filename: str = None):
    """
    Combines multiple Word documents (.docx) into one file, uploads it, and cleans up.
    """
    logger.info(f"Starting DOCX merge task for files: {file_paths}")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Initializing DOCX merge..."})
    
    from docx import Document
    
    app_dir = os.path.dirname(os.path.abspath(__file__))
    backend_root = os.path.dirname(app_dir)
    
    resolved_paths = []
    for fp in file_paths:
        fp = fp.replace('\\', '/')
        if not os.path.isabs(fp):
            fp = os.path.abspath(os.path.join(backend_root, fp))
        resolved_paths.append(fp)
        
    output_filename = f"merged_{uuid.uuid4()}.docx"
    tmp_dir = os.path.join(backend_root, "tmp")
    os.makedirs(tmp_dir, exist_ok=True)
    output_path = os.path.join(tmp_dir, output_filename)
    
    try:
        self.update_state(state="PROGRESS", meta={"progress": 50, "message": "Merging Word files..."})
        
        valid_paths = [fp for fp in resolved_paths if os.path.exists(fp)]
        if not valid_paths:
            raise Exception("No valid Word documents found to merge.")
            
        merged_document = Document(valid_paths[0])
        
        for fp in valid_paths[1:]:
            merged_document.add_page_break()
            sub_doc = Document(fp)
            
            for element in sub_doc.element.body:
                if element.tag.endswith('sectPr'):
                    continue
                merged_document.element.body.append(element)
                
        merged_document.save(output_path)
        
        self.update_state(state="PROGRESS", meta={"progress": 80, "message": "Saving merged Word document..."})
        
        original_name = None
        if original_filename:
            base, ext = os.path.splitext(original_filename)
            original_name = f"{base}_merged{ext}"
            
        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store merged Word document.")
            
        if not is_s3_configured():
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)
            
        return {
            "status": "success",
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"DOCX merge task failed: {e}")
        raise e
    finally:
        for fp in resolved_paths:
            if fp and os.path.exists(fp):
                try:
                    os.remove(fp)
                except Exception:
                    pass
        if os.path.exists(output_path):
            try:
                os.remove(output_path)
            except Exception:
                pass


@celery.task(bind=True)
def docx_to_images_task(self, file_path: str, img_format: str, dpi: int, original_filename: str = None):
    """
    Converts a Word document (.docx) to PDF using LibreOffice, rasterizes PDF pages to images, packages them as ZIP, and uploads.
    """
    logger.info(f"Starting DOCX to image task. Format: {img_format}, DPI: {dpi}")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Converting Word to PDF first..."})
    
    file_path = _resolve_path(file_path)
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")
        
    app_dir = os.path.dirname(os.path.abspath(__file__))
    backend_root = os.path.dirname(app_dir)
    tmp_dir = os.path.join(backend_root, "tmp")
    os.makedirs(tmp_dir, exist_ok=True)
    
    pdf_path = None
    output_filename = ""
    output_path = ""
    
    try:
        pdf_path = _convert_docx_or_pptx_to_pdf(file_path, tmp_dir)
        
        self.update_state(state="PROGRESS", meta={"progress": 60, "message": "Rasterizing PDF pages to images..."})
        
        images = convert_from_path(pdf_path, dpi=dpi)
        
        if len(images) == 1:
            output_filename = f"image_{uuid.uuid4()}.{img_format.lower()}"
            output_path = os.path.join(tmp_dir, output_filename)
            images[0].save(output_path, format=img_format.upper())
            
            original_name = None
            if original_filename:
                base, _ = os.path.splitext(original_filename)
                original_name = f"{base}.{img_format.lower()}"
        else:
            output_filename = f"images_{uuid.uuid4()}.zip"
            output_path = os.path.join(tmp_dir, output_filename)
            with zipfile.ZipFile(output_path, "w") as zipf:
                for idx, img in enumerate(images):
                    img_name = f"page_{idx + 1}.{img_format.lower()}"
                    temp_img_path = os.path.join(tmp_dir, img_name)
                    img.save(temp_img_path, format=img_format.upper())
                    zipf.write(temp_img_path, img_name)
                    os.remove(temp_img_path)
                    
            original_name = None
            if original_filename:
                base, _ = os.path.splitext(original_filename)
                original_name = f"{base}_images.zip"
                
        self.update_state(state="PROGRESS", meta={"progress": 85, "message": "Saving converted images..."})
        
        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store converted images ZIP.")
            
        if not is_s3_configured():
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)
            
        return {
            "status": "success",
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"DOCX to Image conversion failed: {e}")
        raise e
    finally:
        for path in (file_path, pdf_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass


@celery.task(bind=True)
def pptx_to_pdf_task(self, file_path: str, original_filename: str = None):
    """
    Converts a PowerPoint presentation (.pptx) to PDF using LibreOffice, uploads, and cleans up.
    """
    logger.info(f"Starting PPTX to PDF task for {file_path}")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Initializing PPTX to PDF conversion..."})
    
    file_path = _resolve_path(file_path)
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")
        
    app_dir = os.path.dirname(os.path.abspath(__file__))
    backend_root = os.path.dirname(app_dir)
    tmp_dir = os.path.join(backend_root, "tmp")
    os.makedirs(tmp_dir, exist_ok=True)
    
    pdf_path = None
    output_filename = f"converted_{uuid.uuid4()}.pdf"
    output_path = os.path.join(tmp_dir, output_filename)
    
    try:
        self.update_state(state="PROGRESS", meta={"progress": 60, "message": "Converting presentation layout to PDF..."})
        
        generated_pdf = _convert_docx_or_pptx_to_pdf(file_path, tmp_dir)
        shutil.move(generated_pdf, output_path)
        
        self.update_state(state="PROGRESS", meta={"progress": 85, "message": "Saving converted PDF..."})
        
        original_name = None
        if original_filename:
            base, _ = os.path.splitext(original_filename)
            original_name = f"{base}.pdf"
            
        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store converted PDF file.")
            
        if not is_s3_configured():
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)
            
        return {
            "status": "success",
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"PPTX to PDF conversion failed: {e}")
        raise e
    finally:
        for path in (file_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass


@celery.task(bind=True)
def pptx_to_images_task(self, file_path: str, img_format: str, dpi: int, original_filename: str = None):
    """
    Converts a PowerPoint presentation (.pptx) to PDF first, then rasterizes pages to images, packages, and uploads.
    """
    logger.info(f"Starting PPTX to images task for {file_path}")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Converting PowerPoint to PDF first..."})
    
    file_path = _resolve_path(file_path)
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Input file not found: {file_path}")
        
    app_dir = os.path.dirname(os.path.abspath(__file__))
    backend_root = os.path.dirname(app_dir)
    tmp_dir = os.path.join(backend_root, "tmp")
    os.makedirs(tmp_dir, exist_ok=True)
    
    pdf_path = None
    output_filename = ""
    output_path = ""
    
    try:
        pdf_path = _convert_docx_or_pptx_to_pdf(file_path, tmp_dir)
        
        self.update_state(state="PROGRESS", meta={"progress": 60, "message": "Rasterizing presentation slides to images..."})
        
        images = convert_from_path(pdf_path, dpi=dpi)
        
        if len(images) == 1:
            output_filename = f"slide_{uuid.uuid4()}.{img_format.lower()}"
            output_path = os.path.join(tmp_dir, output_filename)
            images[0].save(output_path, format=img_format.upper())
            
            original_name = None
            if original_filename:
                base, _ = os.path.splitext(original_filename)
                original_name = f"{base}.{img_format.lower()}"
        else:
            output_filename = f"slides_{uuid.uuid4()}.zip"
            output_path = os.path.join(tmp_dir, output_filename)
            with zipfile.ZipFile(output_path, "w") as zipf:
                for idx, img in enumerate(images):
                    img_name = f"slide_{idx + 1}.{img_format.lower()}"
                    temp_img_path = os.path.join(tmp_dir, img_name)
                    img.save(temp_img_path, format=img_format.upper())
                    zipf.write(temp_img_path, img_name)
                    os.remove(temp_img_path)
                    
            original_name = None
            if original_filename:
                base, _ = os.path.splitext(original_filename)
                original_name = f"{base}_slides.zip"
                
        self.update_state(state="PROGRESS", meta={"progress": 85, "message": "Saving slide images..."})
        
        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store converted presentation slides.")
            
        if not is_s3_configured():
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)
            
        return {
            "status": "success",
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"PPTX to images conversion failed: {e}")
        raise e
    finally:
        for path in (file_path, pdf_path, output_path):
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass


@celery.task(bind=True)
def merge_pptx_task(self, file_paths: list[str], original_filename: str = None):
    """
    Combines multiple PowerPoint files (.pptx) into one slide deck, uploads it, and cleans up.
    """
    logger.info(f"Starting PPTX merge task for files: {file_paths}")
    self.update_state(state="PROGRESS", meta={"progress": 20, "message": "Initializing PPTX merge..."})
    
    from pptx import Presentation
    import copy
    
    app_dir = os.path.dirname(os.path.abspath(__file__))
    backend_root = os.path.dirname(app_dir)
    
    resolved_paths = []
    for fp in file_paths:
        fp = fp.replace('\\', '/')
        if not os.path.isabs(fp):
            fp = os.path.abspath(os.path.join(backend_root, fp))
        resolved_paths.append(fp)
        
    output_filename = f"merged_{uuid.uuid4()}.pptx"
    tmp_dir = os.path.join(backend_root, "tmp")
    os.makedirs(tmp_dir, exist_ok=True)
    output_path = os.path.join(tmp_dir, output_filename)
    
    try:
        self.update_state(state="PROGRESS", meta={"progress": 50, "message": "Merging presentation slides..."})
        
        valid_paths = [fp for fp in resolved_paths if os.path.exists(fp)]
        if not valid_paths:
            raise Exception("No valid PowerPoint files found to merge.")
            
        prs1 = Presentation(valid_paths[0])
        
        for fp in valid_paths[1:]:
            prs2 = Presentation(fp)
            
            for slide in prs2.slides:
                blank_layout = prs1.slide_layouts[6]
                new_slide = prs1.slides.add_slide(blank_layout)
                
                for shape in slide.shapes:
                    shape_el = copy.deepcopy(shape.element)
                    new_slide.shapes._spTree.append(shape_el)
                    
        prs1.save(output_path)
        
        self.update_state(state="PROGRESS", meta={"progress": 80, "message": "Saving merged presentation..."})
        
        original_name = None
        if original_filename:
            base, ext = os.path.splitext(original_filename)
            original_name = f"{base}_merged{ext}"
            
        download_url = store_processed_file(output_path, output_filename, original_name=original_name)
        if not download_url:
            raise Exception("Failed to store merged PowerPoint.")
            
        if not is_s3_configured():
            local_dest_path = os.path.abspath(os.path.join(backend_root, "static", "outputs", output_filename))
            delete_local_file_task.apply_async((local_dest_path,), countdown=900)
            
        return {
            "status": "success",
            "download_url": download_url
        }
    except Exception as e:
        logger.error(f"PPTX merge task failed: {e}")
        raise e
    finally:
        for fp in resolved_paths:
            if fp and os.path.exists(fp):
                try:
                    os.remove(fp)
                except Exception:
                    pass
        if os.path.exists(output_path):
            try:
                os.remove(output_path)
            except Exception:
                pass
